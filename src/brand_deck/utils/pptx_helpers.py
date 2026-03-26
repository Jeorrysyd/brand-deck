# pptx_helpers.py — 底层 PPT 操作函数
"""Low-level python-pptx helper functions for slide creation and styling."""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from brand_deck.theme import BrandTheme


# Standard slide dimensions (16:9 widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def create_presentation(theme: BrandTheme) -> Presentation:
    """Create a new presentation with 16:9 widescreen dimensions."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs: Presentation) -> object:
    """Add a blank slide to the presentation."""
    layout = prs.slide_layouts[6]  # Blank layout
    return prs.slides.add_slide(layout)


def add_shape(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    fill_color: RGBColor | None = None,
) -> object:
    """Add a rectangle shape to a slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()  # No border

    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()  # Transparent

    return shape


def add_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    font_name: str = "Arial",
    font_size: int = 18,
    font_color: RGBColor = RGBColor(0x33, 0x33, 0x33),
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
) -> object:
    """Add a text box with styled text to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment

    run = p.runs[0] if p.runs else p.add_run()
    if not p.runs:
        run.text = text

    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold

    return txBox


def add_title(
    slide,
    text: str,
    theme: BrandTheme,
    top: int = Inches(0.8),
    left: int = Inches(0.8),
    width: int = Inches(11.7),
    font_size: int = 36,
    color: RGBColor | None = None,
) -> object:
    """Add a title text box with brand styling."""
    return add_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=Inches(1.0),
        text=text,
        font_name=theme.heading_font,
        font_size=font_size,
        font_color=color or theme.text_dark_rgb,
        bold=True,
        alignment=PP_ALIGN.LEFT,
    )


def add_subtitle(
    slide,
    text: str,
    theme: BrandTheme,
    top: int = Inches(1.9),
    left: int = Inches(0.8),
    width: int = Inches(11.7),
    font_size: int = 18,
    color: RGBColor | None = None,
) -> object:
    """Add a subtitle text box with brand styling."""
    return add_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=Inches(0.6),
        text=text,
        font_name=theme.body_font,
        font_size=font_size,
        font_color=color or theme.text_dark_rgb,
        bold=False,
        alignment=PP_ALIGN.LEFT,
    )


def add_body_text(
    slide,
    text: str,
    theme: BrandTheme,
    top: int = Inches(2.8),
    left: int = Inches(0.8),
    width: int = Inches(11.7),
    height: int = Inches(4.0),
    font_size: int = 16,
    color: RGBColor | None = None,
) -> object:
    """Add body text with brand styling."""
    return add_textbox(
        slide,
        left=left,
        top=top,
        width=width,
        height=height,
        text=text,
        font_name=theme.body_font,
        font_size=font_size,
        font_color=color or theme.text_dark_rgb,
        bold=False,
        alignment=PP_ALIGN.LEFT,
    )


def add_image_safe(
    slide,
    image_path: str,
    left: int,
    top: int,
    width: int | None = None,
    height: int | None = None,
) -> object | None:
    """Add an image to a slide. Returns None if image not found (with placeholder)."""
    from pathlib import Path

    if not Path(image_path).exists():
        # Add placeholder text instead of crashing
        add_textbox(
            slide,
            left=left,
            top=top,
            width=width or Inches(5),
            height=height or Inches(3),
            text=f"[Image not found: {Path(image_path).name}]",
            font_size=14,
            font_color=RGBColor(0x99, 0x99, 0x99),
            alignment=PP_ALIGN.CENTER,
        )
        return None

    kwargs = {"left": left, "top": top}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height

    return slide.shapes.add_picture(image_path, **kwargs)


def set_slide_bg(slide, color: RGBColor) -> None:
    """Set the background color of a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color
