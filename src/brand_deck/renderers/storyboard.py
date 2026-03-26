# renderers/storyboard.py — 分镜表格渲染器
"""Render a storyboard (shot list) slide for video production scripts."""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from brand_deck.theme import BrandTheme
from brand_deck.utils.pptx_helpers import add_blank_slide, add_title, set_slide_bg


def render_storyboard(prs, slide_data: dict, theme: BrandTheme) -> None:
    """Render a storyboard table slide."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, theme.bg_rgb)

    title = slide_data.get("title", "Storyboard")
    shots = slide_data.get("shots", [])

    if title:
        add_title(slide, title, theme, top=Inches(0.4), font_size=28)

    if not shots:
        return

    # Storyboard table: Shot # | Visual | Camera | Dialogue | Duration
    headers = ["#", "Visual / Scene", "Camera", "Dialogue / VO", "Duration"]
    num_cols = len(headers)
    num_rows = 1 + len(shots)  # header + data

    # Limit rows per slide (max ~8 shots per slide for readability)
    max_shots = min(len(shots), 8)
    num_rows = 1 + max_shots

    table_left = Inches(0.5)
    table_top = Inches(1.5)
    table_width = Inches(12.3)
    table_height = Inches(5.5)

    col_widths = [Inches(0.7), Inches(4.0), Inches(2.5), Inches(3.5), Inches(1.6)]

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        table_left, table_top,
        table_width, table_height,
    )
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = theme.heading_font
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = theme.text_light_rgb
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.primary_rgb

    # Data rows
    for row_idx, shot in enumerate(shots[:max_shots], start=1):
        values = [
            str(shot.get("shot_number", row_idx)),
            shot.get("visual", ""),
            shot.get("camera", ""),
            shot.get("dialogue", ""),
            shot.get("duration", ""),
        ]

        for col_idx, val in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT if col_idx > 0 else PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = theme.body_font
                    run.font.size = Pt(9)
                    run.font.color.rgb = theme.text_dark_rgb

            # Alternating rows
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
