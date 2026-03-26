# renderers/table.py — 数据表格页渲染器
"""Render a data table slide with branded styling."""

from __future__ import annotations

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from brand_deck.theme import BrandTheme
from brand_deck.utils.pptx_helpers import add_blank_slide, add_title, set_slide_bg


def render_table(prs, slide_data: dict, theme: BrandTheme) -> None:
    """Render a table slide."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, theme.bg_rgb)

    title = slide_data.get("title", "")
    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])

    if title:
        add_title(slide, title, theme, top=Inches(0.5))

    if not headers and not rows:
        return

    # Calculate table dimensions
    num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    num_rows = (1 if headers else 0) + len(rows)

    table_left = Inches(0.8)
    table_top = Inches(1.8)
    table_width = Inches(11.7)
    table_height = min(Inches(5.2), Inches(0.5 * num_rows))

    # Create table
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        table_left, table_top,
        table_width, table_height,
    )
    table = table_shape.table

    # Style header row
    row_idx = 0
    if headers:
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)

            # Header styling
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = theme.heading_font
                    run.font.size = Pt(12)
                    run.font.bold = True
                    run.font.color.rgb = theme.text_light_rgb

            # Header background
            cell.fill.solid()
            cell.fill.fore_color.rgb = theme.primary_rgb
        row_idx = 1

    # Fill data rows
    for data_row in rows:
        for col_idx, cell_text in enumerate(data_row):
            if col_idx >= num_cols:
                break
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = theme.body_font
                    run.font.size = Pt(11)
                    run.font.color.rgb = theme.text_dark_rgb

            # Alternating row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
            else:
                cell.fill.background()

        row_idx += 1
